import os
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
import openpyxl
import pptx


def extract_text_from_file(file_path):
    """
    Extract text from various file types
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    file_ext = os.path.splitext(file_path)[1].lower()
    
    if file_ext == '.pdf':
        return extract_pdf_text(file_path)
    elif file_ext in ['.doc', '.docx']:
        return extract_docx_text(file_path)
    elif file_ext == '.txt':
        return extract_txt_text(file_path)
    elif file_ext in ['.xls', '.xlsx']:
        return extract_excel_text(file_path)
    elif file_ext in ['.ppt', '.pptx']:
        return extract_pptx_text(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_ext}")

def extract_pdf_text(file_path):
    """Extract text from PDF files"""
    try:
        return extract_pdf_text(file_path)
    except Exception as e:
        raise

def extract_docx_text(file_path):
    """Extract text from DOCX files"""
    try:
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    except Exception as e:
        raise

def extract_txt_text(file_path):
    """Extract text from TXT files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read()
        except Exception as e:
            raise

def extract_excel_text(file_path):
    """Extract text from Excel files"""
    try:
        workbook = openpyxl.load_workbook(file_path, read_only=True)
        text = []
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                text.append("\t".join(row_text))
        return "\n".join(text)
    except Exception as e:
        raise

def extract_pptx_text(file_path):
    """Extract text from PowerPoint files"""
    try:
        prs = pptx.Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        raise